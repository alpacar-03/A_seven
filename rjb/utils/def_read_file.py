# 项目主程序
import json
import os
import sys
import mimetypes

# 确保中文显示正常
import matplotlib.pyplot as plt
plt.rcParams["font.family"] = ["SimHei", "WenQuanYi Micro Hei", "Heiti TC"]

# 文件处理库
import pandas as pd
import pdfplumber
import docx
from pptx import Presentation

def read_text_file(file_path):
    """读取文本文件内容"""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
        return file.read()


def read_markdown_file(file_path):
    """读取Markdown文件内容"""
    return read_text_file(file_path)


def read_csv_file(file_path):
    """读取CSV文件内容"""
    try:
        df = pd.read_csv(file_path)
        return df.to_csv(sep='\t', na_rep='nan')
    except Exception as e:
        print(f"读取CSV文件时出错: {e}")
        return f"[无法解析的CSV文件: {os.path.basename(file_path)}]"


def read_excel_file(file_path):
    """读取Excel文件内容"""
    try:
        # 尝试读取所有表
        excel_file = pd.ExcelFile(file_path)
        sheet_names = excel_file.sheet_names
        
        all_content = []
        for sheet_name in sheet_names:
            df = excel_file.parse(sheet_name)
            sheet_content = df.to_csv(sep='\t', na_rep='nan')
            all_content.append(f"工作表 '{sheet_name}':\n{sheet_content}")
        
        return "\n\n".join(all_content)
    except Exception as e:
        print(f"读取Excel文件时出错: {e}")
        return f"[无法解析的Excel文件: {os.path.basename(file_path)}]"


def read_docx_file(file_path):
    """读取DOCX文件内容"""
    try:
        doc = docx.Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        print(f"读取DOCX文件时出错: {e}")
        return f"[无法解析的DOCX文件: {os.path.basename(file_path)}]"


def read_pdf_file(file_path):
    """读取PDF文件内容"""
    try:
        with pdfplumber.open(file_path) as pdf:
            pages = []
            for page in pdf.pages:
                pages.append(page.extract_text())
            return '\n'.join(pages)
    except Exception as e:
        print(f"读取PDF文件时出错: {e}")
        return f"[无法解析的PDF文件: {os.path.basename(file_path)}]"


def read_ppt_file(file_path):
    """读取PPTX文件内容"""
    try:
        prs = Presentation(file_path)
        slides_content = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = [f"幻灯片 {i+1}:"]
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            slides_content.append('\n'.join(slide_text))
            
        return '\n\n'.join(slides_content)
    except Exception as e:
        print(f"读取PPTX文件时出错: {e}")
        return f"[无法解析的PPT文件: {os.path.basename(file_path)}]"


def read_file_content(file_path):
    """根据文件类型读取文件内容"""
    try:
        # 检查文件是否存在
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")
            
        # 获取文件扩展名
        file_ext = os.path.splitext(file_path)[1].lower()
        
        # 根据文件类型选择读取方法
        if file_ext == '.txt':
            return read_text_file(file_path)
        elif file_ext == '.md':
            return read_markdown_file(file_path)
        elif file_ext == '.csv':
            return read_csv_file(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            return read_excel_file(file_path)
        elif file_ext == '.docx':
            return read_docx_file(file_path)
        elif file_ext == '.pdf':
            return read_pdf_file(file_path)
        elif file_ext == '.pptx':
            return read_ppt_file(file_path)
        else:
            # 尝试作为文本文件读取
            return read_text_file(file_path)
                
    except Exception as e:
        print(f"读取文件时出错: {e}")
        return None
